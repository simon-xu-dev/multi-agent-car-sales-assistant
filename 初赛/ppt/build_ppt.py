# -*- coding: utf-8 -*-
"""
初赛方案 PPT 生成脚本（路径 2：python-pptx 可编辑版）
风格：清爽专业风（浅色背景、专业蓝 + 青色强调、结构化卡片/流程布局）
用法：python3 build_ppt.py
输出：基于多Agent的汽车销售自主成交智能助手-初赛方案.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------- 风格常量（清爽专业风） ----------------
PAGE_W, PAGE_H = 13.333, 7.5
C_PAGE = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_SOFT = RGBColor(0xF8, 0xFA, 0xFC)
C_CARD = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_SOFT = RGBColor(0xF8, 0xFA, 0xFC)
C_LINE = RGBColor(0xE2, 0xE8, 0xF0)
C_INK = RGBColor(0x1F, 0x29, 0x37)
C_MUTED = RGBColor(0x47, 0x55, 0x69)
C_PRIMARY = RGBColor(0x25, 0x63, 0xEB)
C_PRIMARY_DK = RGBColor(0x1E, 0x40, 0xAF)
C_PRIMARY_SOFT = RGBColor(0xEF, 0xF6, 0xFF)
C_TEAL = RGBColor(0x0F, 0x76, 0x6E)
C_TEAL_SOFT = RGBColor(0xEC, 0xFD, 0xF5)
C_AMBER = RGBColor(0xB4, 0x53, 0x09)
C_AMBER_SOFT = RGBColor(0xFF, 0xFB, 0xEB)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(PAGE_W)
prs.slide_height = Inches(PAGE_H)
BLANK = prs.slide_layouts[6]


# ---------------- 基础工具 ----------------
def _set_run(run, text, size, color, bold=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn('a:ea'))
    if ea is None:
        ea = rpr.makeelement(qn('a:ea'), {})
        rpr.append(ea)
    ea.set('typeface', FONT)


def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    return box, tf


def para(tf, runs, align=PP_ALIGN.LEFT, first=False, space_before=0.0, space_after=0.0, line=1.0):
    """runs: list of (text, size, color, bold)"""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    for (t, s, c, b) in runs:
        _set_run(p.add_run(), t, s, c, b)
    return p


def rect(slide, x, y, w, h, fill=C_CARD_SOFT, line=C_LINE, lw=0.75,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07, shadow=False):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    sp.shadow.inherit = False
    return sp


def chip(slide, x, y, w, h, text, fill, tcolor, size=11, bold=True, line=None):
    rect(slide, x, y, w, h, fill=fill, line=line, radius=0.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _, tf = tb(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(text, size, tcolor, bold)], align=PP_ALIGN.CENTER, first=True)


def shape_text(sp, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, size=None):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    first = True
    for r in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if isinstance(r, tuple):
            runs_ = [r]
        else:
            runs_ = r
        for (t, s, c, b) in runs_:
            _set_run(p.add_run(), t, s, c, b)


def new_slide(bg=C_PAGE):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, PAGE_W, PAGE_H, fill=bg, line=None, shape=MSO_SHAPE.RECTANGLE)
    return s


FOOTER_TXT = "SalesFlow · 基于多 Agent 的汽车销售自主成交智能助手 · 初赛方案"

def footer(slide, idx):
    rect(slide, 0.55, 7.08, PAGE_W - 1.1, 0.012, fill=C_LINE, line=None, shape=MSO_SHAPE.RECTANGLE)
    _, tf = tb(slide, 0.55, 7.14, 8.5, 0.3)
    para(tf, [(FOOTER_TXT, 9, RGBColor(0x94, 0xA3, 0xB8), False)], first=True)
    _, tf2 = tb(slide, PAGE_W - 1.55, 7.14, 1.0, 0.3)
    para(tf2, [(f"{idx:02d} / 13", 9, RGBColor(0x94, 0xA3, 0xB8), False)],
         align=PP_ALIGN.RIGHT, first=True)


def title_bar(slide, title, tag, subtitle=None):
    # 标题
    _, tf = tb(slide, 0.55, 0.38, 10.6, 0.55)
    para(tf, [(title, 23, C_INK, True)], first=True)
    # 右上标签
    rect(slide, PAGE_W - 2.35, 0.42, 1.8, 0.42, fill=C_PRIMARY_SOFT, line=None, radius=0.5)
    _, tf2 = tb(slide, PAGE_W - 2.35, 0.42, 1.8, 0.42, anchor=MSO_ANCHOR.MIDDLE)
    para(tf2, [(tag, 11, C_PRIMARY, True)], align=PP_ALIGN.CENTER, first=True)
    # 标题下划线
    rect(slide, 0.57, 1.02, 1.35, 0.045, fill=C_PRIMARY, line=None, shape=MSO_SHAPE.RECTANGLE)
    rect(slide, 1.98, 1.02, 0.28, 0.045, fill=C_TEAL, line=None, shape=MSO_SHAPE.RECTANGLE)
    if subtitle:
        _, tf3 = tb(slide, 0.55, 1.14, 12.2, 0.35)
        para(tf3, [(subtitle, 11.5, C_MUTED, False)], first=True)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def card(slide, x, y, w, h, head, head_color, body, head_size=13, body_size=10.5,
         fill=C_CARD, line=C_LINE, num=None, icon=None, body_color=C_MUTED):
    rect(slide, x, y, w, h, fill=fill, line=line, lw=1.0)
    # 顶部色条
    rect(slide, x + 0.18, y + 0.16, 0.09, 0.3, fill=head_color, line=None, shape=MSO_SHAPE.RECTANGLE)
    hx = x + 0.36
    if num:
        _, tfn = tb(slide, hx, y + 0.13, 0.42, 0.34)
        para(tfn, [(num, 15, head_color, True)], first=True)
        hx += 0.44
    _, tf = tb(slide, hx, y + 0.14, w - (hx - x) - 0.18, 0.36)
    para(tf, [(head, head_size, C_INK, True)], first=True)
    _, tfb = tb(slide, x + 0.24, y + 0.56, w - 0.44, h - 0.66)
    para(tfb, [(body, body_size, body_color, False)], first=True, line=1.18)


# ================= Slide 1 封面 =================
s = new_slide()
# 左侧色带装饰
rect(s, 0, 0, 0.18, PAGE_H, fill=C_PRIMARY, line=None, shape=MSO_SHAPE.RECTANGLE)
rect(s, 0.18, 0, 0.07, PAGE_H, fill=C_TEAL, line=None, shape=MSO_SHAPE.RECTANGLE)
# 赛道标签
rect(s, 0.9, 1.35, 3.6, 0.46, fill=C_PRIMARY_SOFT, line=None, radius=0.5)
_, tf = tb(s, 0.9, 1.35, 3.6, 0.46, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Agent Infra 赛道 · 初赛方案", 12, C_PRIMARY, True)], align=PP_ALIGN.CENTER, first=True)
# 主标题
_, tf = tb(s, 0.9, 2.05, 11.5, 1.7)
para(tf, [("基于多 Agent 的", 34, C_INK, True)], first=True, line=1.12)
para(tf, [("汽车销售自主成交智能助手", 40, C_PRIMARY_DK, True)], line=1.12)
# 分隔线
rect(s, 0.92, 4.05, 2.2, 0.05, fill=C_TEAL, line=None, shape=MSO_SHAPE.RECTANGLE)
# 副标题
_, tf = tb(s, 0.9, 4.3, 11.3, 0.8)
para(tf, [("覆盖「线索获取 — 需求分析 — 成交促进 — 售后运营 — 知识沉淀」的自主销售闭环",
           15, C_MUTED, False)], first=True, line=1.3)
# 关键词标签
tags = ["Multi-Agent", "Skill", "MCP", "RAG", "可观测"]
tx = 0.9
for t in tags:
    w = 1.35 if len(t) > 4 else 1.05
    chip(s, tx, 5.25, w, 0.44, t, C_BG_SOFT, C_PRIMARY, size=12, line=C_LINE)
    tx += w + 0.22
# 底部信息
_, tf = tb(s, 0.9, 6.55, 8.0, 0.4)
para(tf, [("SalesFlow 团队  ·  2026", 12, RGBColor(0x94, 0xA3, 0xB8), False)], first=True)
add_notes(s, "各位评委好，我们的作品是「基于多 Agent 的汽车销售自主成交智能助手」。"
             "它面向汽车销售全流程，构建覆盖线索获取、需求分析、成交促进、售后运营与知识沉淀的多 Agent 自主销售闭环，"
             "技术体系由 Multi-Agent、Skill、MCP、RAG 与可观测五大支柱组成。")

# ================= Slide 2 行业痛点 =================
s = new_slide()
title_bar(s, "行业痛点：传统汽车销售的效率困局", "PROBLEM",
          "传统汽车销售高度依赖人工经验，客户咨询分散在官网、微信、电话、短视频平台、门店等多个渠道")
pains = [
    ("线索分散难统一", C_PRIMARY,
     "咨询分散在官网、微信、电话、短视频、门店等渠道，线索无法统一管理，跟进易遗漏。"),
    ("客户画像靠人工", C_TEAL,
     "需求识别与画像构建依赖销售个人经验，判断慢、口径不一，难以形成统一客户视图。"),
    ("销售流程割裂", C_AMBER,
     "车型推荐、优惠报价、试驾预约、金融方案各自为政，缺少一体化推进与状态跟踪。"),
    ("优秀经验难沉淀", C_PRIMARY_DK,
     "成功案例与金牌话术留在个人手里，无法复用与传承，新人培养周期长。"),
]
cx = 0.55
for i, (h, c, b) in enumerate(pains):
    card(s, cx, 1.75, 2.95, 2.1, h, c, b, num=f"0{i+1}")
    cx += 3.11
# 结果条
rect(s, 0.55, 4.35, 12.23, 1.0, fill=C_AMBER_SOFT, line=RGBColor(0xF3, 0xD9, 0xB0), lw=1.0)
_, tf = tb(s, 0.55, 4.35, 12.23, 1.0, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("业务后果：", 15, C_AMBER, True),
          ("线索流失率高    ·    成交周期长    ·    销售效率低", 15, C_AMBER, True)],
     align=PP_ALIGN.CENTER, first=True)
# 量化基线条（对标现状 → Demo 目标）
metrics = [
    ("线索响应", "人工跟进小时级", "多 Agent 分钟级聚合分级"),
    ("经验复用率", "金牌话术≈0%（留在个人手里）", "成交案例自动入库可检索"),
    ("流程衔接", "推荐/报价/试驾/金融各自为政", "一个 Team 房间端到端推进"),
    ("成交周期", "多环节人工等待串联", "低风险动作自动执行压缩等待"),
]
for i, (t, cur, tgt) in enumerate(metrics):
    x = 0.55 + i * 3.11
    rect(s, x, 5.5, 2.95, 1.42, fill=C_CARD, line=C_LINE, lw=1.0)
    _, tf = tb(s, x + 0.18, 5.6, 2.6, 0.3)
    para(tf, [(t, 11.5, C_PRIMARY, True)], first=True)
    _, tf = tb(s, x + 0.18, 5.92, 2.62, 0.95)
    para(tf, [("现状 " + cur, 9.5, C_MUTED, False)], first=True, line=1.1)
    para(tf, [("目标 " + tgt, 9.5, C_TEAL, True)], line=1.1)
footer(s, 2)
add_notes(s, "传统汽车销售有四大痛点：线索分散难统一、客户画像靠人工、销售流程割裂、优秀经验难沉淀，"
             "最终导致线索流失率高、成交周期长、销售效率低。底部给出四个量化改进方向："
             "线索响应从人工小时级到分钟级聚合分级，经验复用率从 0 到案例自动入库，"
             "流程从各自为政到端到端推进，成交周期通过低风险动作自动执行压缩等待。"
             "（注：现状为行业普遍情况描述，复赛将以真实接入数据沉淀量化基线。）")

# ================= Slide 3 总体方案 =================
s = new_slide()
title_bar(s, "总体方案：Multi-Agent 自主成交闭环", "SOLUTION",
          "围绕汽车销售全生命周期构建自主决策系统：Multi-Agent + Skill + MCP + RAG + 可观测")
# 五段闭环（chevron）
stages = [("线索获取", "多渠道线索\n聚合与分级"), ("需求分析", "画像构建\n意图识别"),
          ("成交促进", "策略生成 · 议价\n试驾 · 金融"), ("售后运营", "交付关怀\n复购 · 转介绍"),
          ("知识沉淀", "案例与话术\n入库回流")]
sx, sw, gap = 0.55, 2.62, 0.12
for i, (t, d) in enumerate(stages):
    sp = rect(s, sx + i * (sw + gap), 1.85, sw, 1.35, fill=C_PRIMARY_SOFT if i % 2 == 0 else C_TEAL_SOFT,
              line=C_LINE, shape=MSO_SHAPE.CHEVRON)
    shape_text(sp, [[(t, 14, C_PRIMARY_DK if i % 2 == 0 else C_TEAL, True)],
                    [(d.replace("\n", " · "), 10, C_MUTED, False)]])
# 回流说明
sp = rect(s, 3.3, 3.42, 6.7, 0.46, fill=C_AMBER_SOFT, line=RGBColor(0xF3, 0xD9, 0xB0), radius=0.5)
shape_text(sp, [("↺  经验回流：成功案例与话术沉淀后反哺策略生成，持续优化决策", 11.5, C_AMBER, True)])
# 五大支柱
_, tf = tb(s, 0.55, 4.25, 6.0, 0.4)
para(tf, [("五大技术支柱", 14, C_INK, True)], first=True)
pillars = [("Multi-Agent", "多职能 Agent 协同决策\n端到端任务闭环", C_PRIMARY, C_PRIMARY_SOFT),
           ("Skill", "销售能力抽象层\n可复用 · 可版本化", C_TEAL, C_TEAL_SOFT),
           ("MCP", "标准协议连接\n企业系统", C_PRIMARY_DK, C_PRIMARY_SOFT),
           ("RAG", "知识与案例检索\n可信决策依据", C_TEAL, C_TEAL_SOFT),
           ("可观测", "全链路轨迹\n持续评估优化", C_AMBER, C_AMBER_SOFT)]
px = 0.55
for (t, d, c, bg) in pillars:
    rect(s, px, 4.7, 2.38, 1.75, fill=bg, line=C_LINE)
    _, tf = tb(s, px + 0.16, 4.88, 2.06, 0.4)
    para(tf, [(t, 15, c, True)], first=True)
    _, tf = tb(s, px + 0.16, 5.32, 2.06, 1.0)
    for j, seg in enumerate(d.split("\n")):
        para(tf, [(seg, 10.5, C_MUTED, False)], first=(j == 0), space_before=2)
    px += 2.5
footer(s, 3)
add_notes(s, "总体方案是一个五段式自主成交闭环：线索获取、需求分析、成交促进、售后运营、知识沉淀，"
             "知识沉淀会回流反哺策略生成。闭环之下是五大技术支柱：Multi-Agent 协同、Skill 能力抽象、"
             "MCP 工具连接、RAG 知识增强和全链路可观测。")

# ================= Slide 4 多 Agent 分工 =================
s = new_slide()
title_bar(s, "多 Agent 分工：8 大职能 Agent + 统一编排", "AGENTS",
          "每个 Agent 具备清晰身份定义与能力边界，通过协作完成端到端销售任务闭环")
agents = [
    ("线索聚合 Agent", "多渠道线索归集、清洗、去重与分级，形成统一线索池。", C_PRIMARY),
    ("客户画像 Agent", "构建客户需求、偏好、预算与决策角色画像。", C_TEAL),
    ("购车意图识别 Agent", "识别购车阶段与关键决策信号，判断跟进优先级。", C_PRIMARY),
    ("销售策略生成 Agent", "制定个性化跟进策略与成交路径。", C_TEAL),
    ("智能议价 Agent", "在授权范围内智能议价与优惠匹配。", C_PRIMARY),
    ("订单执行 Agent", "订单、合同、交付流程自动推进与跟踪。", C_TEAL),
    ("客户运营 Agent", "售后关怀、复购激活与转介绍运营。", C_PRIMARY),
    ("知识沉淀 Agent", "成功案例与优秀话术自动提炼入库。", C_TEAL),
]
ax, ay = 0.55, 1.62
for i, (h, b, c) in enumerate(agents):
    row, col = divmod(i, 4)
    card(s, ax + col * 3.11, ay + row * 1.6, 2.95, 1.44, h, c, b, head_size=12.5, body_size=10)
# 拆分理由条（主动回应“是否过度设计”质询）
rect(s, 0.55, 4.94, 12.23, 0.52, fill=C_AMBER_SOFT, line=RGBColor(0xF3, 0xD9, 0xB0))
_, tf = tb(s, 0.8, 4.94, 11.8, 0.52, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("为什么拆成 8 个 Agent：", 11, C_AMBER, True),
          ("单一职责 + 风险边界隔离——议价 Agent 被注入话术也没有订单确认权；每个 Agent 独立承载附录 A 8 字段，可单独评估、灰度与替换。",
           10.5, C_AMBER, False)], first=True, line=1.1)
# AgentTeams 分层编排条
rect(s, 0.55, 5.58, 12.23, 1.16, fill=C_PRIMARY_SOFT, line=C_PRIMARY, lw=1.2)
_, tf = tb(s, 0.8, 5.66, 11.8, 1.0)
para(tf, [("编排基点：AgentTeams（分层架构 Manager → TeamLeader → Worker）", 12, C_PRIMARY_DK, True)],
     first=True, space_after=3)
para(tf, [("Manager 创建 8 个职能 Worker 与独立 TeamLeader Worker（salesflow-demo-leader），在 Team 房间接收任务；"
           "TeamLeader 拆解调度，Worker 专注执行，Human-in-the-loop 随时可介入。",
           10.5, C_MUTED, False)], line=1.2)
footer(s, 4)
add_notes(s, "系统设计了 8 个职能 Agent，覆盖从线索到成交再到知识沉淀的全链路。"
             "拆分不是过度设计，而是两个工程理由：一是单一职责，二是风险边界隔离——"
             "比如议价 Agent 即使被话术注入，也不具备订单确认权限，高风险动作始终在独立 Agent 与审批门禁之后。"
             "编排映射到 AgentTeams 分层架构：Manager 创建 Worker 与 TeamLeader，"
             "TeamLeader 在 Team 房间拆解任务并调度 Worker，人可随时 @mention 介入。")

# ================= Slide 5 端到端闭环 =================
s = new_slide()
title_bar(s, "Demo 场景时间线：DEAL-2001 家庭 SUV 成交之旅", "WORKFLOW",
          "以 family_suv_deal 场景展示任务级自主闭环（3 场景全部通过 selfcheck 36/36 断言）")
timeline = [
    ("① 任务输入", "lead-intake", "官网+企微+电话 3 渠道归并为 1 条线索", "L0"),
    ("② 画像与意图", "profile-builder / intent-analyst", "二胎家庭 · 预算 25-28 万 · 置信度+证据引用", "L0"),
    ("③ 策略生成", "strategy-planner", "理想 L7 / 问界 M7 对比矩阵 · 库存核验 · 报价", "L0-L1"),
    ("④ 自动执行", "negotiation-executor", "试驾预约自动成功 · 授权内优惠自动应用", "L1"),
    ("⑤ 审批门禁", "negotiation-guard", "超授权优惠停止放行 → discount_override 审批任务", "L2"),
    ("⑥ 安全订单", "order-executor", "幂等键创建订单草稿 · 驳回可回滚 · check_deal 验证", "L2"),
    ("⑦ 经验沉淀", "knowledge-miner", "成交案例脱敏入库 · RAG 可检索 · Trace ≥10 条留痕", "L0"),
]
ty = 1.62
for i, (t, ag, d, lv) in enumerate(timeline):
    chip(s, 0.55, ty, 0.85, 0.6, lv, C_TEAL_SOFT if lv in ("L0", "L1") else C_AMBER_SOFT,
         C_TEAL if lv in ("L0", "L1") else C_AMBER, size=10.5, line=C_LINE)
    rect(s, 1.55, ty, 11.23, 0.6, fill=C_CARD, line=C_LINE, lw=1.0)
    rect(s, 1.55, ty, 0.07, 0.6, fill=C_PRIMARY if i % 2 == 0 else C_TEAL, line=None, shape=MSO_SHAPE.RECTANGLE)
    _, tf = tb(s, 1.78, ty + 0.04, 2.55, 0.52, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(t, 11.5, C_INK, True)], first=True)
    _, tf = tb(s, 4.4, ty + 0.04, 5.55, 0.52, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(d, 10.5, C_MUTED, False)], first=True, line=1.05)
    _, tf = tb(s, 10.05, ty + 0.04, 2.6, 0.52, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(ag, 9.5, C_PRIMARY, True)], first=True, line=1.0)
    ty += 0.72
footer(s, 5)
add_notes(s, "这一页用 DEAL-2001（家庭 SUV）场景时间线展示任务级自主闭环的七个关键动作，"
             "每步标注执行 Agent 与风险等级：L0 只读、L1 低风险自动执行、L2 审批门禁。"
             "低风险动作自动执行，超授权优惠立即停止并转审批，订单幂等可回滚，成交案例脱敏沉淀为 RAG 证据。"
             "另外两个场景（首购金融 DEAL-2002、置换复购 DEAL-2003）同样通过全部断言，"
             "三场景合计 36/36，详见代码包 docs/EVIDENCE.md 与 docs/selfcheck_report.txt。")

# ================= Slide 6 AgentTeams 框架能力映射 =================
s = new_slide()
title_bar(s, "AgentTeams 协同设计基点：框架能力映射", "AGENTTEAMS",
          "角色编排、任务拆解、上下文传递、协同执行与状态追踪全部落到 AgentTeams 框架能力")
map_rows = [
    ("角色编排", "Manager 通过 create_agent 串行创建 8 个职能 Worker + 独立 TeamLeader Worker（salesflow-demo-leader），create_team 组建 salesflow-demo Team，对应框架 Manager → TeamLeader → Worker 分层语义"),
    ("任务拆解", "用户在 Team 房间 @team_leader_name 提交销售任务；TeamLeader 按任务特性拆解并调度对应 Worker（按需调用），对应框架 TeamLeader 的指令分解与任务分配"),
    ("上下文传递", "Agent 间经 Matrix Room 事件流传递任务上下文与中间结论（画像/分级/策略/报价/审批状态），对应框架 Matrix 协作架构 + TeamHarness 的 Team/Room/Task 上下文语义"),
    ("协同执行", "Worker 运行时（QwenPaw）调用各自 Skill 与工具网关执行，低风险自动、高风险审批，对应框架 TeamHarness 统一协作协议（Event/Status、File/Tool 引用）"),
    ("状态追踪", "线索/订单/审批状态经工具网关写入 Trace 与 actions，TeamLeader 汇总闭环报告，对应框架 Event/Status 任务事件与状态 + Nacos AI Registry 实例状态回写（Desired → Applied → Status）"),
    ("人机协同", "高风险动作（超授权优惠 L2 / 征信 L2 / 订单确认 L2 / 合同 L3）人工审批与回滚，任何时刻可 @mention 介入，对应框架 Human-in-the-loop 设计"),
]
rows = len(map_rows) + 1
tbl_shape = s.shapes.add_table(rows, 2, Inches(0.55), Inches(1.62), Inches(12.23), Inches(4.9))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.2)
tbl.columns[1].width = Inches(10.03)
for j, htxt in enumerate(["协同环节", "AgentTeams 框架能力映射"]):
    c = tbl.cell(0, j)
    c.fill.solid(); c.fill.fore_color.rgb = C_PRIMARY
    tfc = c.text_frame; tfc.word_wrap = True
    p = tfc.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), htxt, 12, C_WHITE, True)
for i, (a, b) in enumerate(map_rows):
    for j, val in enumerate((a, b)):
        c = tbl.cell(i + 1, j)
        c.fill.solid()
        c.fill.fore_color.rgb = C_CARD_SOFT if i % 2 else C_CARD
        tfc = c.text_frame; tfc.word_wrap = True
        p = tfc.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
        _set_run(p.add_run(), val, 10.5, C_INK if j == 0 else C_MUTED, j == 0)
footer(s, 6)
add_notes(s, "评审红线：AgentTeams 不是提名字就行，必须看到角色编排、任务拆解、上下文传递怎么落到框架能力上。"
             "本页给出六维映射：角色编排=create_agent/create_team 分层；任务拆解=TeamLeader 指令分解；"
             "上下文传递=Matrix Room + TeamHarness；协同执行=Worker 运行时 + 工具网关；"
             "状态追踪=Event/Status + Nacos 状态回写；人机协同=Human-in-the-loop 审批与介入。")

# ================= Slide 7 Skill 体系 =================
s = new_slide()
title_bar(s, "Skill 体系：可复用的销售能力抽象层", "SKILLS",
          "Skill 作为任务能力抽象层：明确输入输出、调用条件、依赖工具、失败处理与安全边界")
skills = [
    ("lead-fusion", "线索归并", "多渠道会话合并去重分级", "输入：会话列表\n输出：合并线索", "缺客户 ID 启发式归并转人工"),
    ("profile-building", "画像构建", "结构化画像 + 置信度", "输入：会话 + 历史\n输出：画像 + 证据", "历史为空置信度上限 0.6"),
    ("intent-scoring", "意图识别", "信号字典打分与分级", "输入：画像 + 信号\n输出：意向度 + 优先级", "信号不足降级 nurture"),
    ("car-recommendation", "车型推荐", "按画像匹配候选车型", "输入：画像 + 库存\n输出：对比矩阵 + 理由", "知识不足不编造参数"),
    ("quote-pricing", "报价优惠", "政策内报价与优惠审批", "输入：政策 + 客户等级\n输出：报价单", "超授权自动转 L2 审批"),
    ("negotiation-guard", "议价风控", "让步轨迹与底线守护", "输入：报价 + 已让步\n输出：让步决策", "触底即停转人工"),
    ("test-drive-booking", "试驾预约", "档期查询预约与回滚", "输入：客户 + 门店\n输出：预约单", "失败重试推荐替代档期"),
    ("finance-plan", "金融方案", "方案对比与征信审批", "输入：成交价 + 期数\n输出：方案对比", "未授权不发起征信"),
    ("order-safe-execute", "订单执行", "幂等 + 审批门禁 + 回滚", "输入：报价 + 确认\n输出：订单状态", "审批驳回回滚草稿"),
    ("case-mining", "案例沉淀", "闭环复盘脱敏入库", "输入：报告 + Trace\n输出：结构化案例", "证据不足仅沉淀摘要"),
    ("deal-memory", "记忆检索", "案例 + 客户记忆 RAG", "输入：查询 + 客户 ID\n输出：Top-N 证据", "检索失败降级标注"),
]
rows = len(skills) + 1
tbl_shape = s.shapes.add_table(rows, 5, Inches(0.55), Inches(1.58), Inches(12.23), Inches(4.85))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(1.85)
tbl.columns[1].width = Inches(1.85)
tbl.columns[2].width = Inches(2.6)
tbl.columns[3].width = Inches(3.2)
tbl.columns[4].width = Inches(2.73)
headers = ["Skill 名称", "能力", "用途", "输入 / 输出", "失败处理"]
for j, htxt in enumerate(headers):
    c = tbl.cell(0, j)
    c.fill.solid(); c.fill.fore_color.rgb = C_PRIMARY
    tfc = c.text_frame; tfc.word_wrap = True
    p = tfc.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), htxt, 11.5, C_WHITE, True)
for i, row in enumerate(skills):
    for j, val in enumerate(row):
        c = tbl.cell(i + 1, j)
        c.fill.solid()
        c.fill.fore_color.rgb = C_CARD_SOFT if i % 2 else C_CARD
        tfc = c.text_frame; tfc.word_wrap = True
        first = True
        for seg in val.split("\n"):
            p = tfc.paragraphs[0] if first else tfc.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.LEFT if j > 1 else PP_ALIGN.CENTER
            _set_run(p.add_run(), seg, 9.5, C_INK if j == 0 else C_MUTED, j == 0)
_, tf = tb(s, 0.55, 6.46, 12.2, 0.6)
para(tf, [("双层 Skill 结构：", 10.5, C_INK, True),
          ("基础云能力复用阿里云官方用云 Skills（模型调用 / 短信触达 / 证据归档）不重复造轮子；业务层自研 11 个 Skill 聚焦销售场景。",
           10.5, C_MUTED, False)], first=True, line=1.1)
para(tf, [("复用价值：每个 Skill 按 10 字段模板定义（含调用条件 / 依赖工具 / 权限与安全 / 复用价值），可跨 Agent、跨场景复用，"
           "支持版本化与灰度发布（Nacos AI Registry）。", 10.5, C_TEAL, True)], line=1.1)
footer(s, 7)
add_notes(s, "Skill 是本方案的必选能力层。生态上我们是双层结构：基础云能力直接复用阿里云官方用云 Skills，"
             "不重复构建；业务层把线索归并、画像构建、意图识别、车型推荐、报价优惠、议价风控、试驾预约、"
             "金融方案、订单执行、案例沉淀、记忆检索 11 个关键能力沉淀为自研可复用 Skill，"
             "每个都按 10 字段模板定义，并通过 Nacos AI Registry 做版本管理与灰度发布。")

# ================= Slide 8 MCP 工具集成 =================
s = new_slide()
title_bar(s, "MCP 工具集成：统一连接企业系统", "MCP",
          "Skill 承担任务能力抽象层，MCP 承担工具连接层，Higress 承担统一网关治理")
# 左：能力层
rect(s, 0.55, 1.85, 3.1, 3.4, fill=C_PRIMARY_SOFT, line=C_PRIMARY, lw=1.2)
_, tf = tb(s, 0.75, 2.0, 2.7, 3.1)
para(tf, [("Agent / Skill 层", 14, C_PRIMARY_DK, True)], first=True, space_after=6)
for t in ["能力抽象与任务编排", "按 Skill 契约调用工具", "不感知具体系统实现"]:
    para(tf, [("· " + t, 11, C_MUTED, False)], space_before=4)
# 箭头1
sp = rect(s, 3.78, 3.3, 0.62, 0.5, fill=C_PRIMARY, line=None, shape=MSO_SHAPE.RIGHT_ARROW)
# 中：网关 + 协议层
rect(s, 4.52, 1.85, 4.3, 3.4, fill=C_CARD, line=C_PRIMARY, lw=1.5)
_, tf = tb(s, 4.72, 2.0, 3.9, 3.1)
para(tf, [("Higress AI 网关 + MCP", 13.5, C_PRIMARY_DK, True)], first=True, space_after=6)
for t in ["统一入口 · 鉴权 · 路由", "限流（防工具滥用）", "参数 Schema / 返回结构", "审计日志 · 失败重试 · 幂等控制"]:
    para(tf, [("· " + t, 11, C_MUTED, False)], space_before=4)
# 箭头2
sp = rect(s, 8.95, 3.3, 0.62, 0.5, fill=C_TEAL, line=None, shape=MSO_SHAPE.RIGHT_ARROW)
# 右：企业系统
rect(s, 9.7, 1.85, 3.08, 3.4, fill=C_TEAL_SOFT, line=C_TEAL, lw=1.2)
_, tf = tb(s, 9.9, 2.0, 2.7, 0.4)
para(tf, [("企业系统", 14, C_TEAL, True)], first=True)
systems = ["CRM 客户关系", "库存 / DMS", "金融审批", "保险", "合同系统", "企业微信"]
gx, gy = 9.9, 2.5
for i, t in enumerate(systems):
    r, cidx = divmod(i, 2)
    rect(s, gx + cidx * 1.42, gy + r * 0.62, 1.3, 0.5, fill=C_CARD, line=C_LINE, radius=0.18)
    _, tf = tb(s, gx + cidx * 1.42, gy + r * 0.62, 1.3, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(t, 9.5, C_INK, False)], align=PP_ALIGN.CENTER, first=True)
# 底部说明
rect(s, 0.55, 5.6, 12.23, 1.05, fill=C_BG_SOFT, line=C_LINE)
_, tf = tb(s, 0.8, 5.6, 11.8, 1.05, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("接口契约标准化：", 12.5, C_INK, True),
          ("工具定义名称、调用入口、参数 Schema、返回结构、权限范围与审计记录；工具调用链按「工具.函数」命名空间组织并预置 future_mcp_mapping，"
           "迁移 MCP Server 仅需协议适配，业务 Skill 与 Agent 零改动。",
           12, C_MUTED, False)], first=True, line=1.25)
footer(s, 8)
add_notes(s, "MCP 是工具连接层：Agent 和 Skill 通过 MCP 标准协议统一连接 CRM、库存 DMS、金融审批、"
             "保险、合同与企业微信。Higress 作为统一 AI 网关提供鉴权、路由、限流与观测，"
             "协议层内置 Schema、审计、重试、幂等与降级策略，工具调用链按「工具.函数」命名空间组织，"
             "迁移 MCP Server 仅需协议适配，不重构调用链。")

# ================= Slide 9 RAG 知识增强 =================
s = new_slide()
title_bar(s, "RAG 知识增强：为销售决策提供可信依据", "RAG",
          "RAG 作为 Agent、Skill、MCP 调用链中的上下文能力，回答与决策只依据检索证据")
sources = [
    ("产品知识", "车型 / 配置 / 竞品对比", C_PRIMARY),
    ("销售 SOP", "流程规范 / 标准话术", C_TEAL),
    ("历史成交案例", "成功路径 / 议价经验", C_PRIMARY_DK),
    ("补贴与政策", "优惠 / 金融 / 置换政策", C_AMBER),
]
for i, (t, d, c) in enumerate(sources):
    x = 0.55 + i * 3.11
    rect(s, x, 1.75, 2.95, 1.5, fill=C_CARD, line=C_LINE, lw=1.0)
    rect(s, x, 1.75, 2.95, 0.09, fill=c, line=None, shape=MSO_SHAPE.RECTANGLE)
    _, tf = tb(s, x + 0.2, 2.0, 2.55, 1.1)
    para(tf, [(t, 13.5, C_INK, True)], first=True, space_after=4)
    para(tf, [(d, 11, C_MUTED, False)])
# 机制链
mech = [("MCP 接入数据源", "连接知识库与\n业务数据"), ("Skill 封装检索", "证据对齐\n与结果写入"),
        ("Agent 判断决策", "评估证据是否\n足以支撑结论")]
for i, (t, d) in enumerate(mech):
    x = 1.15 + i * 4.05
    sp = rect(s, x, 3.95, 3.05, 1.15, fill=C_PRIMARY_SOFT if i < 2 else C_TEAL_SOFT, line=C_LINE,
              shape=MSO_SHAPE.CHEVRON)
    shape_text(sp, [[(t, 13, C_PRIMARY_DK if i < 2 else C_TEAL, True)],
                    [(d.replace("\n", "，"), 10, C_MUTED, False)]])
    if i < 2:
        pass
# 证据边界
rect(s, 0.55, 5.55, 12.23, 1.05, fill=C_AMBER_SOFT, line=RGBColor(0xF3, 0xD9, 0xB0))
_, tf = tb(s, 0.8, 5.55, 11.8, 1.05, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("证据边界：", 12.5, C_AMBER, True),
          ("回答与决策只依据检索到的证据；未被真实来源支持的事实（价格、库存、政策）一律不编造；检索失败时明确降级而不是猜测。",
           12, C_AMBER, False)], first=True, line=1.25)
footer(s, 9)
add_notes(s, "RAG 覆盖四类知识源：产品知识、销售 SOP、历史成交案例和补贴政策。"
             "机制上由 MCP 接入数据源，Skill 封装检索与证据对齐，Agent 判断证据是否足以支撑决策，"
             "并坚持证据边界：没有证据支持的事实一律不编造。方案覆盖 RAG 4 项能力中的"
             "「Agent 记忆存储」（客户历史时间窗口+语义检索）与「知识库 RAG」（案例向量化检索），"
             "存储层可由 PolarDB 承载，支持时间窗口查询与向量语义检索。")

# ================= Slide 10 全链路可观测 =================
s = new_slide()
title_bar(s, "全链路可观测：持续评估与优化", "OBSERVABILITY",
          "覆盖 Agent 推理、Skill 调用、工具执行与 RAG 检索的全链路推理轨迹")
objs = [("Agent 推理轨迹", "意图识别 / 策略生成\n的推理过程留痕", C_PRIMARY),
        ("Skill 调用", "调用条件 · 输入输出\n成功率与耗时", C_TEAL),
        ("MCP 工具执行", "企业系统调用审计\n异常与降级记录", C_PRIMARY_DK),
        ("RAG 检索", "检索命中与证据引用\n回答可溯源", C_AMBER)]
for i, (t, d, c) in enumerate(objs):
    x = 0.55 + i * 3.11
    rect(s, x, 1.7, 2.95, 1.6, fill=C_CARD, line=C_LINE, lw=1.0)
    rect(s, x + 0.18, 1.88, 0.09, 0.3, fill=c, line=None, shape=MSO_SHAPE.RECTANGLE)
    _, tf = tb(s, x + 0.36, 1.86, 2.4, 0.36)
    para(tf, [(t, 12.5, C_INK, True)], first=True)
    _, tf = tb(s, x + 0.24, 2.32, 2.5, 0.9)
    for j, seg in enumerate(d.split("\n")):
        para(tf, [(seg, 10.5, C_MUTED, False)], first=(j == 0), space_before=2)
# 数据类型 + 应用
rect(s, 0.55, 3.62, 5.95, 2.5, fill=C_BG_SOFT, line=C_LINE)
_, tf = tb(s, 0.8, 3.8, 5.5, 2.2)
para(tf, [("数据类型（OTel GenAI 语义）", 13, C_INK, True)], first=True, space_after=6)
for t in ["Trace：Agent / Skill / MCP / RAG / LLM Span", "Log：决策依据 / 失败原因 / 权限与审批事件（TraceId 关联）",
          "Metrics：会话数 / 时延 / Token 消耗 / 工具成功率"]:
    para(tf, [("· " + t, 10.5, C_MUTED, False)], space_before=5)
rect(s, 6.83, 3.62, 5.95, 2.5, fill=C_TEAL_SOFT, line=C_LINE)
_, tf = tb(s, 7.08, 3.8, 5.5, 2.2)
para(tf, [("应用场景", 13, C_TEAL, True)], first=True, space_after=6)
for t in ["在线监控与告警：异常链路实时发现", "离线评估：LLM-as-Judge / 规则评估量化效果",
          "评估结果回流 Dataset，驱动 Prompt / Skill 迭代"]:
    para(tf, [("· " + t, 10.5, C_MUTED, False)], space_before=5)
# 工具落地条
rect(s, 0.55, 6.3, 12.23, 0.72, fill=C_PRIMARY_SOFT, line=C_PRIMARY, lw=1.0)
_, tf = tb(s, 0.8, 6.3, 11.8, 0.72, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("落地路线：", 11.5, C_PRIMARY_DK, True),
          ("开源自建 LoongSuite + AgentScope Studio（OTel 自动埋点 + 可视化 Trace / Agent 评估）；生产环境平滑迁移阿里云 AgentLoop（Agent-as-a-Judge 持续优化）。",
           10.5, C_MUTED, False)], first=True, line=1.2)
footer(s, 10)
add_notes(s, "可观测覆盖四类对象：Agent 推理轨迹、Skill 调用、MCP 工具执行和 RAG 检索，"
             "沉淀 Trace、Log、Metrics 三类数据（遵循 OpenTelemetry GenAI 语义），"
             "支撑在线监控告警、离线量化评估（LLM-as-Judge / 规则评估）与 badcase 复盘迭代。"
             "技术路线：开源采用 LoongSuite + AgentScope Studio，生产可迁移阿里云 AgentLoop。")

# ================= Slide 11 推荐工具链治理 =================
s = new_slide()
title_bar(s, "推荐工具链：全链路治理架构", "TOOLCHAIN",
          "按需利用推荐工具链，重点在治理设计与接口关系，而非堆叠数量")
chains = [
    ("Nacos AI Registry", "AI 管理中心", C_PRIMARY,
     ["Agent / Skill / Prompt / MCP 统一注册发现", "多版本灰度 · 秒级回滚 · 权限隔离", "案例经审核发布为 Skill 资产"]),
    ("Higress AI 网关", "统一入口", C_TEAL,
     ["模型 / Agent / 工具调用统一入口", "鉴权 · 路由 · 限流 · 观测", "工具级权限边界与密钥隔离"]),
    ("PolarDB 数据层", "存储底座", C_PRIMARY_DK,
     ["向量 / 长记忆 / RAG / 审计日志存储", "案例与记忆向量化 + 相似度检索", "接口契约化，可替换性明确"]),
    ("RocketMQ 消息队列", "事件驱动", C_AMBER,
     ["审批事件 / 订单状态流转 / 触达任务", "可靠投递（重试 / 死信）+ 状态幂等", "执行状态与 TraceId 关联可追踪"]),
]
for i, (t, d, c, items) in enumerate(chains):
    x = 0.55 + i * 3.11
    rect(s, x, 1.75, 2.95, 3.25, fill=C_CARD, line=C_LINE, lw=1.0)
    rect(s, x, 1.75, 2.95, 0.62, fill=c, line=None, shape=MSO_SHAPE.RECTANGLE)
    _, tf = tb(s, x, 1.75, 2.95, 0.62, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(t, 12, C_WHITE, True)], align=PP_ALIGN.CENTER, first=True)
    para(tf, [(d, 9.5, C_WHITE, False)], align=PP_ALIGN.CENTER)
    _, tf = tb(s, x + 0.2, 2.55, 2.6, 2.3)
    for j, it in enumerate(items):
        para(tf, [("· " + it, 9.5, C_MUTED, False)], first=(j == 0), space_before=6, line=1.15)
# 底部：云 Skills 与工具观
rect(s, 0.55, 5.3, 12.23, 1.4, fill=C_BG_SOFT, line=C_LINE)
_, tf = tb(s, 0.8, 5.38, 11.8, 1.25)
para(tf, [("阿里云官方用云 Skills：", 11, C_INK, True),
          ("复用官方 Skills 承载基础云能力（模型调用 / 短信触达 / 证据归档），业务 Skill 聚焦销售场景；鉴权走 Higress 统一入口，凭证隔离。",
           10.5, C_MUTED, False)], first=True, space_after=4)
para(tf, [("工具观：", 11, C_INK, True),
          ("每个工具说明必要性、接口契约、可替换性、权限边界与闭环证据，后续迁移成本清晰可评估。",
           10.5, C_MUTED, False)], first=True)
footer(s, 11)
add_notes(s, "推荐工具链按需利用：Nacos AI Registry 做 Agent/Skill/Prompt 治理与灰度发布，"
             "Higress 做统一网关（鉴权/路由/限流/观测），PolarDB 做向量与审计存储，"
             "RocketMQ 做事件驱动与状态流转，可观测采用 LoongSuite + AgentScope Studio 或 AgentLoop，"
             "阿里云官方用云 Skills 承载基础云能力。评审重点是设计理念、接口契约、可替换性与权限边界，而非堆叠数量。")

# ================= Slide 12 创新与复用 =================
s = new_slide()
title_bar(s, "创新点与开放复用价值", "INNOVATION",
          "区别于传统智能客服仅回答问题，本方案打造能够自主推进成交的 AI 销售助手")
rows = [
    ("定位", "问答机器人，被动回答", "自主推进成交的 AI 销售助手"),
    ("能力边界", "咨询问答", "策略生成 / 智能议价 / 试驾预约 / 订单执行"),
    ("经验积累", "依赖人工总结", "成功案例自动沉淀，回流优化决策"),
    ("技术架构", "单体机器人", "Multi-Agent + Skill + MCP + RAG + 可观测"),
]
tbl_shape = s.shapes.add_table(5, 3, Inches(0.55), Inches(1.62), Inches(12.23), Inches(2.7))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(1.9)
tbl.columns[1].width = Inches(4.4)
tbl.columns[2].width = Inches(5.93)
for j, htxt in enumerate(["维度", "传统智能客服", "本方案"]):
    c = tbl.cell(0, j)
    c.fill.solid(); c.fill.fore_color.rgb = C_PRIMARY if j != 2 else C_TEAL
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), htxt, 12, C_WHITE, True)
for i, (a, b, cc) in enumerate(rows):
    vals = [a, b, cc]
    for j, val in enumerate(vals):
        c = tbl.cell(i + 1, j)
        c.fill.solid()
        c.fill.fore_color.rgb = C_TEAL_SOFT if j == 2 else (C_CARD_SOFT if i % 2 else C_CARD)
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
        _set_run(p.add_run(), val, 11, C_INK if j != 1 else C_MUTED, j in (0, 2))
# 开放复用价值
_, tf = tb(s, 0.55, 4.62, 6.0, 0.4)
para(tf, [("开放 / 复用价值", 14, C_INK, True)], first=True)
vals = [
    ("模块化设计", "Skill 能力可独立复用，MCP 接口遵循标准化协议，可快速接入 CRM / ERP / DMS / 金融 / 保险系统。", C_PRIMARY),
    ("跨行业迁移", "闭环范式可迁移至房地产、保险、金融产品、高端零售等复杂销售场景。", C_TEAL),
    ("开源计划", "开源 Skill 清单、MCP 接口契约与示例场景，作为企业智能销售平台基础能力持续扩展。", C_PRIMARY_DK),
]
for i, (t, d, c) in enumerate(vals):
    card(s, 0.55 + i * 4.15, 5.05, 3.98, 1.62, t, c, d, head_size=12.5, body_size=10)
footer(s, 12)
add_notes(s, "创新点在于：这不是问答机器人，而是能自主推进成交的 AI 销售助手，"
             "并且把 Skill、MCP、RAG、可观测融合成能力抽象—工具连接—知识增强—持续优化的完整技术闭环。"
             "模块化与标准化设计让它可快速接入企业系统，并迁移到房地产、保险、高端零售等场景。")

# ================= Slide 13 进展与路线图 =================
s = new_slide()
title_bar(s, "初赛实现边界与复赛计划", "ROADMAP",
          "初赛 V0.2：可信方案设计 + 可复现自检证据；复赛 V0.5：真实接入与运行验证")
# 左栏：初赛已实现
rect(s, 0.55, 1.75, 6.05, 3.15, fill=C_TEAL_SOFT, line=C_LINE, lw=1.0)
rect(s, 0.55, 1.75, 6.05, 0.56, fill=C_TEAL, line=None, shape=MSO_SHAPE.RECTANGLE)
_, tf = tb(s, 0.55, 1.75, 6.05, 0.56, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("初赛已实现 · V0.2（证据可复现）", 13, C_WHITE, True)], align=PP_ALIGN.CENTER, first=True)
done = [
    "8 Agent（附录 A 8 字段）+ 11 Skill（附录 B 10 字段）+ AgentTeams 编排材料齐备",
    "HTTP mock 工具网关：9 类工具 · 3 场景 · MCP 等价契约与迁移路径",
    "3 场景端到端业务闭环：审批 / 回滚 / 议价底线 / 案例入库",
    "离线自检 36/36 断言通过，运行证据落盘（docs/EVIDENCE.md）",
    "AgentTeams 实际运行：3 场景 79 次工具调用（DEAL-2001:33 / DEAL-2002:24 / DEAL-2003:22）",
]
_, tf = tb(s, 0.83, 2.5, 5.5, 2.3)
for j, it in enumerate(done):
    para(tf, [("✓ " + it, 11, C_MUTED, False)], first=(j == 0), space_before=8, line=1.2)
# 右栏：复赛计划
rect(s, 6.73, 1.75, 6.05, 3.15, fill=C_PRIMARY_SOFT, line=C_LINE, lw=1.0)
rect(s, 6.73, 1.75, 6.05, 0.56, fill=C_PRIMARY, line=None, shape=MSO_SHAPE.RECTANGLE)
_, tf = tb(s, 6.73, 1.75, 6.05, 0.56, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("复赛计划 · V0.5（8.24 晋级后）", 13, C_WHITE, True)], align=PP_ALIGN.CENTER, first=True)
plan = [
    "MCP Server 替换 mock 适配器：真实 CRM / 库存 / 金融审批接入",
    "AgentTeams 实际部署：Manager 创建 Team 现场跑通 3 场景",
    "全链路可观测：LoongSuite / AgentScope Studio + 评估回流",
    "量化指标沉淀：线索转化率 / 成交周期 / 工具成功率",
]
_, tf = tb(s, 7.01, 2.5, 5.5, 2.3)
for j, it in enumerate(plan):
    para(tf, [("→ " + it, 11, C_MUTED, False)], first=(j == 0), space_before=8, line=1.2)
# 证据索引条
rect(s, 0.55, 5.02, 12.23, 0.66, fill=C_CARD, line=C_LINE)
_, tf = tb(s, 0.8, 5.02, 11.8, 0.66, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("证据索引（docs/EVIDENCE.md）：", 11, C_PRIMARY, True),
          ("自检 36/36 断言 · AgentTeams 实际运行 3 场景 79 次工具调用 · 每场景 Trace ≥10 条 · L2 审批 / 订单回滚 / 底线转人工用例各在环 · 3 场景期望信号表可逐项对照",
           10.5, C_MUTED, False)], first=True, line=1.1)
# 边界声明条
rect(s, 0.55, 5.8, 12.23, 0.9, fill=C_AMBER_SOFT, line=RGBColor(0xF3, 0xD9, 0xB0))
_, tf = tb(s, 0.8, 5.8, 11.8, 0.9, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("边界声明：", 11.5, C_AMBER, True),
          ("初赛聚焦 mock 环境的任务级自主闭环（方案设计 + 可复现断言）；会话级自由对话、真实系统接入与生产级风控在复赛推进。",
           11, C_AMBER, False)], first=True, line=1.2)
footer(s, 13)
add_notes(s, "初赛已完成：mock 工具网关（9 类工具、3 场景）、8 Agent + 11 Skill 完整材料、"
             "3 场景端到端业务闭环（审批/回滚/议价底线/案例入库），36/36 断言全部通过并落盘为运行证据；"
             "证据索引见代码包 docs/EVIDENCE.md，含每场景 Trace ≥10 条、L2 审批、回滚与底线转人工用例。"
             "复赛计划：MCP 真实系统接入、AgentTeams 实际部署演示、全链路可观测与评估回流、量化指标沉淀。"
             "我们如实声明边界：初赛是 mock 环境的任务级自主闭环，会话级自由对话与生产级风控留给复赛。"
             "让 AI 不止于回答，而是自主推进每一次成交。谢谢！")

# ---------------- 保存 ----------------
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "基于多Agent的汽车销售自主成交智能助手-初赛方案.pptx")
prs.save(OUT)
print("已生成:", OUT)
print("页数:", len(prs.slides.__iter__.__self__._sldIdLst))
