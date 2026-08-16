#!/usr/bin/env python3
# 修复初赛 PPT 过时内容：
# 1) P13 运行证据数据刷新为 2026-08-16 三场景全闭环（116 次调用/24 节点/4 类故障自愈）
# 2) P13 删除与现状矛盾的"复赛才 LLM 驱动"表述
# 3) P4/P6 旧容器名 salesflow-demo → carsales-demo
# 4) P9 三处文案错字
import sys
from pptx import Presentation

PPTX = "/Users/chery-not-23982/Learn/competation/Agent-infra/SalesFlow/ppt和作品简介/基于多Agent的汽车销售自主成交智能助手-初赛方案.pptx"

REPLACEMENTS = [
    # --- P13 核心刷新：旧数据 → 8/16 全闭环新数据 ---
    ("AgentTeams 真框架已跑（11 容器 + 127 条 transcript）",
     "AgentTeams 真框架 3 场景一天全闭环（11 容器 v1.1.2 · 8/16 单日）"),
    ("✓ AgentTeams 真框架跑通：3 场景 · 11 容器 · Matrix 协议 127 条 transcript",
     "✓ AgentTeams 真框架跑通：3 场景全闭环 · 24 DAG 节点全绿 · 116 次工具调用 100% 成功 · 4 类真实故障代码级自愈"),
    ("AgentTeams 真框架 3 场景 127 条 transcript",
     "AgentTeams 真框架 3 场景全闭环（116 次工具调用 100% 成功）"),
    # --- P13 删除与现状矛盾的复赛项（Worker 已是真 LLM ReAct 驱动） ---
    ("→ Worker 脚本调度替换为 LLM 推理（全自主闭环 · TeamLeader ReAct → Worker 自主决策）",
     "→ 多模型路由与调度策略优化 · Agent/Skill/LLM 层可观测 Span 补齐"),
    # --- P4/P6 容器名刷新 ---
    ("salesflow-demo-leader", "carsales-demo-leader"),
    ("create_team 组建 salesflow-demo Team", "create_team 组建 carsales-demo Team"),
    # --- P9 文案错字 ---
    ("连接知识库与，业务数据", "连接知识库与业务数据"),
    ("证据对齐，与结果写入", "证据对齐与结果写入"),
    ("评估证据是否，足以支撑结论", "评估证据是否足以支撑结论"),
]


def iter_paragraphs(shape):
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            yield p
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    yield p


def replace_in_paragraph(p, old, new):
    full = ''.join(r.text for r in p.runs)
    if old not in full:
        return False
    new_full = full.replace(old, new)
    if not p.runs:
        return False
    p.runs[0].text = new_full
    for r in p.runs[1:]:
        r.text = ''
    return True


def main():
    prs = Presentation(PPTX)
    hits = {old: 0 for old, _ in REPLACEMENTS}
    for slide in prs.slides:
        for shape in slide.shapes:
            for p in iter_paragraphs(shape):
                for old, new in REPLACEMENTS:
                    if replace_in_paragraph(p, old, new):
                        hits[old] += 1
    prs.save(PPTX)
    print("=== 替换结果 ===")
    for old, n in hits.items():
        status = "OK" if n > 0 else "!! 未命中"
        print(f"[{status}] x{n}  {old[:40]}")
    print("SAVED:", PPTX)


if __name__ == "__main__":
    main()
