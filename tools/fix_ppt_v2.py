#!/usr/bin/env python3
# 按参赛手册官方案例（OpsPilot Zero）叙事范式优化 PPT：
# 1) 每页副标题改为"思路引导句"（先讲为什么/怎么想，再讲怎么做）
# 2) Demo 场景页补客户故事线（仿官方案例"杭州某电商公司的客诉"）
# 3) 去掉初赛边界页的技术术语堆砌，聚焦"可信交付"思路
from pptx import Presentation

PPTX = "/Users/chery-not-23982/Learn/competation/Agent-infra/SalesFlow/ppt和作品简介/基于多Agent的汽车销售自主成交智能助手-初赛方案.pptx"

REPLACEMENTS = [
    # --- P3 总体方案：从技术名词列表 → 闭环设计思路 ---
    ("围绕汽车销售全生命周期构建自主决策系统：Multi-Agent + Skill + MCP + RAG + 可观测",
     "核心思路：把销售闭环拆成五段——自主能执行、协同可验证、经验可沉淀，越卖越聪明"),
    # --- P4 Agent 分工：从身份描述 → 拆分原则 ---
    ("每个 Agent 具备清晰身份定义与能力边界，通过协作完成端到端销售任务闭环",
     "把人工销售链条拆成职责单一、风险边界隔离、可单独评估与替换的协作单元"),
    # --- P5 Demo 场景：补客户故事线（对齐官方案例叙事） ---
    ("以 family_suv_deal 场景展示任务级自主闭环（3 场景全部通过 selfcheck 56/56 断言）",
     "二胎家庭张先生想换一台 25 万级大空间 SUV：从多渠道咨询到成交，系统全程自主推进、关键节点风控把关"),
    # --- P6 AgentTeams 映射：强调真映射而非提名 ---
    ("角色编排、任务拆解、上下文传递、协同执行与状态追踪全部落到 AgentTeams 框架能力",
     "不只提框架名字：编排、拆解、上下文、执行、追踪逐项映射到 AgentTeams 能力，并有真框架运行证据"),
    # --- P7 Skill：从字段罗列 → 沉淀思路 ---
    ("Skill 作为任务能力抽象层：明确输入输出、调用条件、依赖工具、失败处理与安全边界",
     "把专家销售经验变成可复用、可评估、可版本化的工程能力，而不是一次性提示词"),
    # --- P8 MCP：从分层描述 → 契约思路 ---
    ("Skill 承担任务能力抽象层，MCP 承担工具连接层，Higress 承担统一网关治理",
     "用统一接口契约屏蔽企业系统差异：未来迁移 MCP Server 只需协议适配，业务侧零改动"),
    # --- P10 可观测：从覆盖范围 → 治理思路 ---
    ("覆盖 Agent 推理、Skill 调用、工具执行与 RAG 检索的全链路推理轨迹",
     "自主成交的关键不是无限自动化，而是让每一步可验证、可审计、可回放"),
    # --- P13 初赛证据行：去术语堆砌 ---
    ("✓ 自检 56/56 + 编排 G1/G2/G3 + 三个 LLM 决策点(审批三安全分支+车型推荐+工具顺序) + OTel 运行时 39 span + 重放 95 span + OSS REST 签名验证 + RAG/记忆 + 评估 1.0",
     "✓ 可复现证据：自检 56/56 断言 + 三个 LLM 决策点（审批三安全分支 / 车型推荐 / 工具顺序）+ 全链路 Trace 留痕 + 评估得分 1.0"),
    # --- P13 边界声明：去括号堆砌 ---
    ("边界声明：初赛聚焦 mock 环境任务级自主闭环（方案设计+可复现断言+三个 LLM 决策点已驱动+OTel 运行时插桩+OSS REST 签名验证+AgentTeams 真框架已跑）；会话级自由对话、真实系统接入、生产级风控在复赛推进。",
     "边界声明：初赛聚焦 mock 环境任务级自主闭环（真 AgentTeams 框架运行 + 可复现断言）；会话级自由对话、真实系统接入、生产级风控在复赛推进。"),
    # --- P14 安全页副标题：突出可控性思路 ---
    ("高风险动作禁止默认放行——LLM 自主推理 + 决策与执行分离 + append-only 审计轨迹",
     "自主成交的前提是可控：风险分级定边界、审批门禁管高危、决策与执行分离、审计全程留痕"),
]


def iter_paragraphs(shape):
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            yield p
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    yield p


def replace_in_paragraph(p, old, new):
    full = ''.join(r.text for r in p.runs)
    if old not in full:
        return False
    new_full = full.replace(old, new)
    if not p.runs:
        return False
    p.runs[0].text = new_full
    for r in p.runs[1:]:
        r.text = ''
    return True


def main():
    prs = Presentation(PPTX)
    hits = {old: 0 for old, _ in REPLACEMENTS}
    for slide in prs.slides:
        for shape in slide.shapes:
            for p in iter_paragraphs(shape):
                for old, new in REPLACEMENTS:
                    if replace_in_paragraph(p, old, new):
                        hits[old] += 1
    prs.save(PPTX)
    print("=== 叙事优化替换结果 ===")
    for old, n in hits.items():
        status = "OK" if n > 0 else "!! 未命中"
        print(f"[{status}] x{n}  {old[:45]}")
    print("SAVED:", PPTX)


if __name__ == "__main__":
    main()
