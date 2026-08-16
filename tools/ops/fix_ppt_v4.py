#!/usr/bin/env python3
"""fix_ppt_v4.py — 重构 PPT 为 9 页新结构

重构策略：
  从 14 页中保留 9 页，删除 5 页，更新标题/标签/副标题/页码/关键内容。
  保留原有视觉设计（深色背景、圆角矩形、配色方案）。

原始 14 页（0-based index）：
  0:P1封面  1:P2目录  2:P3概览  3:P4 Agent分工
  4:P5 Demo  5:P6 AgentTeams  6:P7 Skill  7:P8 MCP
  8:P9 RAG  9:P10可观测  10:P11工具链  11:P12创新
  12:P13路线图  13:P14安全

删除 5 页：P2(目录) P5(Demo) P9(RAG) P11(工具链) P12(创新)
删除索引（0-based）：1, 4, 8, 10, 11

新 9 页结构：
  P1 封面          ← 原P1 (idx 0)  保持不变
  P2 场景与价值     ← 原P3 (idx 2)  改标题/内容
  P3 整体架构       ← 原P4 (idx 3)  改标题/副标题
  P4 Agent协同      ← 原P6 (idx 5)  改标题/标签
  P5 Skill与工具    ← 原P7 (idx 6)  改标题/标签/底部
  P6 可观测性       ← 原P10 (idx 9) 改页码
  P7 可行性与落地   ← 原P13 (idx 12) 改标题/内容
  P8 安全边界       ← 原P14 (idx 13) 改标题/标签
  P9 致谢           ← 原P14 (idx 13) ... 等等
"""

import shutil
from pptx import Presentation

PPTX = "/Users/chery-not-23982/Learn/competation/Agent-infra/SalesFlow/ppt和作品简介/基于多Agent的汽车销售自主成交智能助手-初赛方案.pptx"
BAK = PPTX + ".bak"

# ── 要删除的幻灯片索引（0-based）──
# 当前 PPT 已经是 9 页（之前的版本已处理过），无需再删除
DELETE_INDICES = []  # [1, 4, 8, 10, 11]  # P2, P5, P9, P11, P12 — 已在前序版本中删除


# ── 工具函数 ──────────────────────────────────────────────

def find_shapes_by_text(slide, text):
    """找到包含指定文本的所有 shape"""
    results = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and text in shape.text:
            results.append(shape)
    return results


def update_text(shape, new_text):
    """更新 shape 的文本，保留首个 run 的格式"""
    if not shape or not shape.has_text_frame:
        return False
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = new_text
        for r in tf.paragraphs[0].runs[1:]:
            r.text = ''
        return True
    return False


def replace_in_all_text(slide, old_text, new_text):
    """在 slide 所有 shape（含表格）中做精确替换"""
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                full = ''.join(r.text for r in p.runs)
                if old_text in full:
                    new_full = full.replace(old_text, new_text)
                    if p.runs:
                        p.runs[0].text = new_full
                        for r in p.runs[1:]:
                            r.text = ''
                    count += 1
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        full = ''.join(r.text for r in p.runs)
                        if old_text in full:
                            new_full = full.replace(old_text, new_text)
                            if p.runs:
                                p.runs[0].text = new_full
                                for r in p.runs[1:]:
                                    r.text = ''
                            count += 1
    return count


# ── 主逻辑 ──────────────────────────────────────────────

def main():
    # 1. 备份
    shutil.copy2(PPTX, BAK)
    print(f"[1/4] 备份已保存: {BAK}")

    prs = Presentation(PPTX)
    slides = list(prs.slides)
    print(f"[2/4] 原始幻灯片: {len(slides)} 页")

    # 2. 删除不需要的幻灯片（降序删除避免索引偏移）
    from lxml import etree
    nsmap = {'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
    sldIdLst = prs._element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
    for idx in sorted(DELETE_INDICES, reverse=True):
        sldId = list(sldIdLst)[idx]
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)
        print(f"  删除原 P{idx + 1}")

    slides = list(prs.slides)
    print(f"  剩余: {len(slides)} 页")

    # 新顺序映射：
    # new[0] = 原P1  (封面)
    # new[1] = 原P3  (概览 → 场景与价值)
    # new[2] = 原P4  (Agent分工 → 整体架构)
    # new[3] = 原P6  (AgentTeams → Agent协同)
    # new[4] = 原P7  (Skill → Skill与工具)
    # new[5] = 原P8  (MCP → 工具集成，内容合并到P5说明)
    # new[6] = 原P10 (可观测)
    # new[7] = 原P13 (路线图 → 可行性)
    # new[8] = 原P14 (安全)

    # 3. 更新各页内容
    print("[3/4] 更新页面内容...")

    # ─── P1 封面（new[0]）─── 保持不变
    print("  ✓ P1 封面（保持不变）")

    # ─── P2 场景与价值（new[1] = 原P3）───
    s = slides[1]
    # 标题：方案概览 → 场景与价值
    for shape in find_shapes_by_text(s, "方案概览"):
        update_text(shape, "场景与价值：汽车销售的 AI 机遇")
        break
    # 标签：OVERVIEW → SCENARIO
    for shape in find_shapes_by_text(s, "OVERVIEW"):
        update_text(shape, "SCENARIO")
        break
    # 副标题
    for shape in find_shapes_by_text(s, "Agent 统一编排"):
        update_text(shape,
                    "线索分散 · 画像靠经验 · 报价不统一 · 知识随人走\n"
                    "面向 4S 店与经销商集团，构建自主成交闭环")
        break
    # 四个能力卡片改为痛点/场景关键词
    card_updates = {
        "Agent": "4 大痛点",
        "Skill": "4 类场景",
        "MCP": "业务价值",
        "RAG": "行业迁移",
    }
    for old_label, new_label in card_updates.items():
        for shape in find_shapes_by_text(s, old_label):
            if len(shape.text.strip()) <= 15:
                update_text(shape, new_label)
                break
    # 卡片描述
    desc_updates = {
        "统一编排": "线索分散 · 画像靠经验\n报价不统一 · 知识随人走",
        "可复用": "新客接待 · 首购金融\n老客置换 · 复合联动",
        "标准协议": "转化率提升 · 成本降低\n风险可控 · 知识沉淀",
        "知识与案例": "房产 · 保险 · 金融\n高端零售 · 教育培训",
    }
    for old_text, new_text in desc_updates.items():
        for shape in find_shapes_by_text(s, old_text):
            if len(shape.text.strip()) < 30:
                update_text(shape, new_text)
                break
    # 底部说明
    for shape in find_shapes_by_text(s, "端到端自主闭环"):
        update_text(shape,
                    "覆盖「线索获取 — 需求分析 — 成交促进 — 售后运营 — 知识沉淀」全链路\n"
                    "行业可迁移：房产、保险、金融产品等复杂销售场景")
        break
    # 页码
    for shape in find_shapes_by_text(s, "03 / 14"):
        update_text(shape, "02 / 09")
        break
    print("  ✓ P2 场景与价值")

    # ─── P3 整体架构（new[2] = 原P4）───
    s = slides[2]
    for shape in find_shapes_by_text(s, "多 Agent 分工"):
        update_text(shape, "方案设计 — 整体架构")
        break
    for shape in find_shapes_by_text(s, "AGENTS"):
        update_text(shape, "ARCHITECTURE")
        break
    for shape in find_shapes_by_text(s, "把人工销售链条"):
        update_text(shape,
                    "8 大职能 Agent 按销售流程分工：线索→画像→意图→策略→议价→订单→运营→知识\n"
                    "AgentTeams 三层编排（Manager→TeamLeader→Worker）· 单一职责 + 风险边界隔离")
        break
    # 更新底部说明
    for shape in find_shapes_by_text(s, "为什么拆成 8 个"):
        tf = shape.text_frame
        if tf.paragraphs:
            p = tf.paragraphs[0]
            if p.runs:
                p.runs[0].text = ("为什么拆成 8 个 Agent：单一职责 + 风险边界隔离——议价 Agent 没有订单确认权；"
                                  "每个 Agent 可单独评估、灰度与替换")
                for r in p.runs[1:]:
                    r.text = ''
            if len(tf.paragraphs) > 1:
                p2 = tf.paragraphs[1]
                if p2.runs:
                    p2.runs[0].text = ("编排基点：AgentTeams 三层架构 · Manager 创建 8 个 Worker + TeamLeader\n"
                                       "上下文传递：经 Matrix Room 事件流传递画像/策略/报价/审批等中间结论")
                    for r in p2.runs[1:]:
                        r.text = ''
        break
    for shape in find_shapes_by_text(s, "04 / 14"):
        update_text(shape, "03 / 09")
        break
    print("  ✓ P3 整体架构")

    # ─── P4 Agent协同（new[3] = 原P6）───
    s = slides[3]
    for shape in find_shapes_by_text(s, "AgentTeams 协同设计基点"):
        update_text(shape, "方案设计 — Agent 协同闭环")
        break
    for shape in find_shapes_by_text(s, "AGENTTEAMS"):
        update_text(shape, "COLLABORATION")
        break
    for shape in find_shapes_by_text(s, "不只提框架名字"):
        update_text(shape,
                    "端到端闭环：线索→画像→意图→策略→议价→订单→运营→知识沉淀\n"
                    "结果验证：87 项断言 · Golden/Badcase 测试 · 异常分支：故障注入 7 用例覆盖降级/重试/隔离/恢复")
        break
    for shape in find_shapes_by_text(s, "06 / 14"):
        update_text(shape, "04 / 09")
        break
    print("  ✓ P4 Agent协同")

    # ─── P5 Skill与工具（new[4] = 原P7）───
    s = slides[4]
    for shape in find_shapes_by_text(s, "Skill 体系"):
        update_text(shape, "Skill 与工具集成：销售能力工程体系")
        break
    for shape in find_shapes_by_text(s, "SKILLS"):
        update_text(shape, "SKILLS & TOOLS")
        break
    for shape in find_shapes_by_text(s, "把专家销售经验"):
        update_text(shape,
                    "13 个 Skill 固化销售 SOP + 25 个工具函数覆盖 9 类企业系统\n"
                    "MCP 等价契约统一连接，迁移 MCP Server 只需协议适配，业务侧零改动")
        break
    # 更新底部说明
    replace_in_all_text(s,
                        "业务层自研 11 个 Skill 聚焦销售场景",
                        "业务层自研 13 个 Skill 固化销售 SOP（含 sms-approval-alert + 官方短信触达 Skill）")
    # 更新底部 MCP/RAG 描述中的工具数
    replace_in_all_text(s, "22 工具(9 类)", "25 工具(9 类)")
    for shape in find_shapes_by_text(s, "07 / 14"):
        update_text(shape, "05 / 09")
        break
    print("  ✓ P5 Skill与工具")

    # ─── P6 MCP 工具集成（new[5] = 原P8）───
    s = slides[5]
    for shape in find_shapes_by_text(s, "MCP 工具集成"):
        update_text(shape, "工具集成：统一连接企业系统")
        break
    for shape in find_shapes_by_text(s, "08 / 14"):
        update_text(shape, "06 / 09")
        break
    print("  ✓ P6 MCP工具集成")

    # ─── P7 可观测性（new[6] = 原P10）───
    s = slides[6]
    # 标题保持不变，更新 span 数据和页码
    replace_in_all_text(s,
                        "OTel 运行时 39 span + 重放 95 span",
                        "OTel 运行时 39 span + 三层 trace 树 261 span（Agent→Skill→Tool，挂载率 100%）")
    for shape in find_shapes_by_text(s, "10 / 14"):
        update_text(shape, "07 / 09")
        break
    print("  ✓ P7 可观测性")

    # ─── P8 可行性与落地（new[7] = 原P13）───
    s = slides[7]
    for shape in find_shapes_by_text(s, "初赛实现边界"):
        update_text(shape, "可行性与落地计划")
        break
    for shape in find_shapes_by_text(s, "ROADMAP"):
        update_text(shape, "FEASIBILITY")
        break
    for shape in find_shapes_by_text(s, "初赛 V0.2"):
        update_text(shape, "当前进展：方案设计 + 可复现证据 + Demo 可运行")
        break
    # 更新证据区域
    for shape in find_shapes_by_text(s, "证据索引"):
        update_text(shape,
                    "✓ 自检 87/87 断言 · 4 场景全闭环（含 DEAL-2004 复合联动）\n"
                    "✓ 116 次工具调用 100% 成功 · 35 项可复现证据\n"
                    "✓ Demo 前端可运行 · L2 审批/回滚全验证")
        break
    # 复赛计划 → 开放/开源计划
    for shape in find_shapes_by_text(s, "复赛计划"):
        update_text(shape, "开放 / 开源计划")
        break
    for shape in find_shapes_by_text(s, "→ MCP Server"):
        update_text(shape,
                    "→ 13 Skill + MCP 契约 + pgvector PoC + 故障注入框架全开放\n"
                    "→ 真实 4S 店 pilot · 社区共建 · 跨行业迁移（房产/保险/金融）")
        break
    for shape in find_shapes_by_text(s, "边界声明"):
        update_text(shape,
                    "后续路线图：真实系统接入 → 生产级风控 → 跨行业复制\n"
                    "初赛聚焦 mock 环境任务级自主闭环（真 AgentTeams 框架 + 可复现断言）")
        break
    for shape in find_shapes_by_text(s, "13 / 14"):
        update_text(shape, "08 / 09")
        break
    print("  ✓ P8 可行性与落地")

    # ─── P9 安全边界（new[8] = 原P14）───
    s = slides[8]
    for shape in find_shapes_by_text(s, "安全闭环"):
        update_text(shape, "安全边界与风控：L0-L3 风险分级 × 审批门禁 × 审计")
        break
    for shape in find_shapes_by_text(s, "SECURITY"):
        update_text(shape, "BOUNDARY")
        break
    for shape in find_shapes_by_text(s, "自主成交的前提"):
        update_text(shape,
                    "风险分级定边界 · 审批门禁管高危 · 决策与执行分离 · 审计全程留痕\n"
                    "L0 自动 → L1 通知 → L2 审批 → L3 转人工 · 议价底线守护 · 超授权自动停止")
        break
    # 更新 span 引用
    replace_in_all_text(s,
                        "OTel 运行时 39 span + 重放 95 span",
                        "OTel 运行时 39 span + 三层 trace 树 261 span")
    for shape in find_shapes_by_text(s, "14 / 14"):
        update_text(shape, "09 / 09")
        break
    print("  ✓ P9 安全边界")

    # 4. 保存
    prs.save(PPTX)
    print(f"\n[4/4] 已保存: {PPTX}")

    # 5. 验证
    prs2 = Presentation(PPTX)
    print(f"\n=== V4 重构结果 ===")
    print(f"总页数: {len(prs2.slides)}")
    for si, slide in enumerate(prs2.slides, 1):
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                t = shape.text.strip().split('\n')[0][:60]
                if len(t) > 5 and "SalesFlow" not in t and not t.startswith("0"):
                    title = t
                    break
        print(f"  P{si}: {title}")


if __name__ == "__main__":
    main()
