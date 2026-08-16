#!/usr/bin/env python3
# 提交前最终数据一致性修复（对齐仓库可核验事实）：
# 1) 工具数：25 → 33 工具函数 / 11 类系统（tool_catalog.json 实测）
# 2) Skill 数：11 → 13（skills/ 实测）
# 3) 真框架场景数：4 → 3（8/16 真跑证据仅 3 场景；DEAL-2004 仅自检/故障注入级覆盖）
from pptx import Presentation

PPTX = "/Users/chery-not-23982/Learn/competation/Agent-infra/SalesFlow/ppt和作品简介/基于多Agent的汽车销售自主成交智能助手-初赛方案.pptx"

REPLACEMENTS = [
    ("13 个 Skill 固化销售 SOP + 25 个工具函数覆盖 9 类企业系统",
     "13 个 Skill 固化销售 SOP + 33 个工具函数覆盖 11 类企业系统"),
    ("✓ 8 Agent（附录 A 8 字段）+ 11 Skill（附录 B 10 字段）+ AgentTeams 真框架 4 场景一天全闭环（11 容器 v1.1.2 · 8/16 单日）",
     "✓ 8 Agent（附录 A 8 字段）+ 13 Skill（附录 B 10 字段）+ AgentTeams 真框架 3 场景一天全闭环（11 容器 v1.1.2 · 8/16 单日）"),
    ("✓ HTTP mock 工具网关：25 工具(9 类) · 4 场景（含 DEAL-2004 复合联动） · MCP 等价契约 + 三个 LLM 决策点(审批门禁+车型推荐+工具调用顺序)",
     "✓ HTTP mock 工具网关：33 工具函数(11 类系统) · 4 场景覆盖（含 DEAL-2004 复合联动） · MCP 等价契约 + 三个 LLM 决策点(审批门禁+车型推荐+工具调用顺序)"),
    ("✓ 4 场景端到端业务闭环：审批 / 回滚 / 议价底线 / 案例入库",
     "✓ 3 场景端到端业务闭环：审批 / 回滚 / 议价底线 / 案例入库（DEAL-2004 复合联动另经自检与故障注入覆盖）"),
    ("✓ AgentTeams 真框架跑通：4 场景全闭环 · 24 DAG 节点全绿 · 116 次工具调用 100% 成功 · 4 类真实故障代码级自愈",
     "✓ AgentTeams 真框架跑通：3 场景全闭环 · 24 DAG 节点全绿 · 116 次工具调用 100% 成功 · 4 类真实故障代码级自愈"),
    ("✓ 自检 87/87 断言 · 4 场景全闭环（含 DEAL-2004 复合联动）",
     "✓ 自检 87/87 断言 · 4 场景覆盖（含 DEAL-2004 复合联动）"),
]


def iter_paragraphs(shape):
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            yield p
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    yield p


def replace_in_paragraph(p, old, new):
    full = ''.join(r.text for r in p.runs)
    if old not in full:
        return False
    new_full = full.replace(old, new)
    if not p.runs:
        return False
    p.runs[0].text = new_full
    for r in p.runs[1:]:
        r.text = ''
    return True


def main():
    prs = Presentation(PPTX)
    hits = {old: 0 for old, _ in REPLACEMENTS}
    for slide in prs.slides:
        for shape in slide.shapes:
            for p in iter_paragraphs(shape):
                for old, new in REPLACEMENTS:
                    if replace_in_paragraph(p, old, new):
                        hits[old] += 1
    prs.save(PPTX)
    print("=== 数据一致性修复 ===")
    for old, n in hits.items():
        status = "OK" if n > 0 else "!! 未命中"
        print(f"[{status}] x{n}  {old[:48]}")
    print("SAVED:", PPTX)


if __name__ == "__main__":
    main()
