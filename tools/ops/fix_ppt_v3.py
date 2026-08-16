#!/usr/bin/env python3
"""fix_ppt_v3.py — 增量刷新 PPT 中 6 处过时数据点（+ 可选增强）

数据点更新清单：
  P7  Skill 数量 11 → 13（新增 sms-approval-alert + 1 个官方 Skill）
  P10 可观测 span 95 → 261（三层 trace 树：Agent→Skill→Tool，挂载率 100%）
  P13 场景数 3 → 4（新增 DEAL-2004 复合联动）
  P13 断言数 56 → 87
  P13 工具数 22 → 25
  P14 安全页 span 引用 95 → 261（保持一致）

可选增强：
  P13 补充 35 项可复现证据、Demo 前端可运行
"""
import shutil
from pptx import Presentation

PPTX = "/Users/chery-not-23982/Learn/competation/Agent-infra/SalesFlow/ppt和作品简介/基于多Agent的汽车销售自主成交智能助手-初赛方案.pptx"
BAK  = PPTX + ".bak"

# ── 精确替换对：(旧文本, 新文本) ──────────────────────────────
REPLACEMENTS = [
    # ─── P7 Skill 数量 ───
    ("业务层自研 11 个 Skill 聚焦销售场景",
     "业务层自研 13 个 Skill 聚焦销售场景（含 sms-approval-alert + 官方短信触达 Skill）"),

    # ─── P10 可观测 span ───
    ("OTel 运行时 39 span + 重放 95 span（GenAI semconv）",
     "OTel 运行时 39 span + 三层 trace 树 261 span（Agent→Skill→Tool，挂载率 100%，GenAI semconv）"),

    # ─── P13 场景 / 断言 / 工具 ───
    # 证据索引行（长模式优先，避免被短模式抢先）
    ("自检 56/56 断言 · AgentTeams 真框架 3 场景全闭环（116 次工具调用 100% 成功）",
     "自检 87/87 断言 · AgentTeams 真框架 4 场景全闭环（含 DEAL-2004 复合联动 · 116 次工具调用 100% 成功）"),
    # 场景 3→4
    ("AgentTeams 真框架 3 场景一天全闭环",
     "AgentTeams 真框架 4 场景一天全闭环"),
    ("22 工具(9 类) · 3 场景",
     "25 工具(9 类) · 4 场景（含 DEAL-2004 复合联动）"),
    ("3 场景端到端业务闭环",
     "4 场景端到端业务闭环"),
    # 断言 56→87（可复现证据行）
    ("自检 56/56 断言 + 三个 LLM 决策点",
     "自检 87/87 断言 + 三个 LLM 决策点"),
    # 工具 22→25
    ("HTTP mock 工具网关：22 工具(9 类)",
     "HTTP mock 工具网关：25 工具(9 类)"),
    # 可选增强：补充 35 项证据 + Demo
    ("每场景 Trace ≥10 条",
     "35 项可复现证据 · 每场景 Trace ≥10 条 · Demo 前端可运行"),
    # P13 AgentTeams 跑通行
    ("AgentTeams 真框架跑通：3 场景全闭环",
     "AgentTeams 真框架跑通：4 场景全闭环"),

    # ─── P14 安全页 span 引用（保持一致）───
    ("OTel 运行时 39 span + 重放 95 span",
     "OTel 运行时 39 span + 三层 trace 树 261 span"),
]


def iter_paragraphs(shape):
    """递归产出 shape 中所有段落（含表格单元格）。"""
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            yield p
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    yield p


def replace_in_paragraph(p, old, new):
    """在段落的 runs 拼接文本中做精确替换。"""
    full = ''.join(r.text for r in p.runs)
    if old not in full:
        return False
    new_full = full.replace(old, new)
    if not p.runs:
        return False
    # 保留首 run 格式，合并文本到首 run
    p.runs[0].text = new_full
    for r in p.runs[1:]:
        r.text = ''
    return True


def main():
    # 备份
    shutil.copy2(PPTX, BAK)
    print(f"备份已保存: {BAK}")

    prs = Presentation(PPTX)
    hits = {old: 0 for old, _ in REPLACEMENTS}

    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            for p in iter_paragraphs(shape):
                for old, new in REPLACEMENTS:
                    if replace_in_paragraph(p, old, new):
                        hits[old] += 1

    prs.save(PPTX)

    # 打印结果
    print("\n=== V3 数据点刷新结果 ===")
    all_ok = True
    for old, n in hits.items():
        status = "✓" if n > 0 else "✗ 未命中"
        if n == 0:
            all_ok = False
        print(f"  [{status}] ×{n}  {old[:55]}")
    print(f"\nSAVED: {PPTX}")
    if all_ok:
        print("ALL REPLACEMENTS HIT ✓")
    else:
        print("WARNING: some replacements missed — check text above")


if __name__ == "__main__":
    main()
